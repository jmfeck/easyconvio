from __future__ import annotations

import os
import subprocess
import shutil
from typing import Optional, List

from .base import BaseFile


def _require_command(cmd: str, install_hint: str) -> None:
    if shutil.which(cmd) is None:
        raise RuntimeError(f"'{cmd}' not found. {install_hint}")


def _libreoffice_binary() -> str:
    """Return the LibreOffice CLI name, preferring 'libreoffice' over 'soffice'.

    Windows installs ship `soffice.exe` instead of `libreoffice`.
    """
    for name in ("libreoffice", "soffice"):
        if shutil.which(name):
            return name
    raise RuntimeError(
        "LibreOffice not found. Install LibreOffice: https://www.libreoffice.org/"
    )


class PresentationFile(BaseFile):
    """Presentation file with manipulation and conversion methods."""

    def _load(self) -> None:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "PresentationFile requires python-pptx. "
                "Install with: pip install easyconvio[presentations]"
            )
        self._tmp_pptx: Optional[str] = None
        # python-pptx natively reads pptx and ppsx (same OOXML container)
        if self.format in ("pptx", "ppsx"):
            self._prs = Presentation(self.path)
        else:
            # ppt, pps, odp — convert to pptx via LibreOffice first
            self._tmp_pptx = self._libreoffice_to_pptx()
            self._prs = Presentation(self._tmp_pptx)

    def _libreoffice_to_pptx(self) -> str:
        import tempfile
        binary = _libreoffice_binary()
        out_dir = tempfile.mkdtemp()
        subprocess.run(
            [binary, "--headless", "--convert-to", "pptx",
             "--outdir", out_dir, self.path],
            check=True,
            capture_output=True,
        )
        base = os.path.splitext(os.path.basename(self.path))[0]
        return os.path.join(out_dir, f"{base}.pptx")

    def close(self) -> None:
        if self._tmp_pptx and os.path.exists(self._tmp_pptx):
            try:
                shutil.rmtree(os.path.dirname(self._tmp_pptx))
            except OSError:
                pass

    # --- Properties ---

    @property
    def slide_count(self) -> int:
        """Number of slides."""
        return len(self._prs.slides)

    @property
    def slide_width(self) -> int:
        """Slide width in EMU."""
        return self._prs.slide_width

    @property
    def slide_height(self) -> int:
        """Slide height in EMU."""
        return self._prs.slide_height

    # --- Operations ---

    def extract_text(self) -> List[str]:
        """Extract text from each slide as a list of strings."""
        texts = []
        for slide in self._prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text.append(paragraph.text)
            texts.append("\n".join(slide_text))
        return texts

    def extract_images(self, output_dir: str = ".") -> List[str]:
        """Extract all images from the presentation to a directory."""
        os.makedirs(output_dir, exist_ok=True)
        paths = []
        img_count = 0
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    img = shape.image
                    ext = img.content_type.split("/")[-1]
                    img_count += 1
                    filename = f"image_{img_count}.{ext}"
                    filepath = os.path.join(output_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(img.blob)
                    paths.append(filepath)
        return paths

    def remove_slide(self, index: int) -> PresentationFile:
        """Remove a slide by index."""
        rId = self._prs.slides._sldIdLst[index].rId
        self._prs.part.drop_rel(rId)
        del self._prs.slides._sldIdLst[index]
        return self

    # --- Export ---

    def to_pptx(self, output_path: Optional[str] = None) -> str:
        """Export as PPTX."""
        output_path = self._output_path("pptx", output_path)
        self._prs.save(output_path)
        return output_path

    def _libreoffice_convert(self, fmt: str, output_path: Optional[str] = None) -> str:
        binary = _libreoffice_binary()
        output_path = self._output_path(fmt, output_path)
        out_dir = os.path.dirname(os.path.abspath(output_path))
        subprocess.run(
            [binary, "--headless", "--convert-to", fmt, "--outdir", out_dir, self.path],
            check=True,
            capture_output=True,
        )
        base = os.path.splitext(os.path.basename(self.path))[0]
        generated = os.path.join(out_dir, f"{base}.{fmt}")
        if generated != os.path.abspath(output_path):
            os.rename(generated, output_path)
        return output_path

    def to_pdf(self, output_path: Optional[str] = None) -> str:
        """Export as PDF (requires LibreOffice)."""
        return self._libreoffice_convert("pdf", output_path)

    def to_odp(self, output_path: Optional[str] = None) -> str:
        """Export as ODP (requires LibreOffice)."""
        return self._libreoffice_convert("odp", output_path)

    def to_ppt(self, output_path: Optional[str] = None) -> str:
        """Export as PPT (requires LibreOffice)."""
        return self._libreoffice_convert("ppt", output_path)

    def to_html(self, output_path: Optional[str] = None) -> str:
        """Export as HTML (requires LibreOffice)."""
        return self._libreoffice_convert("html", output_path)

    def to_pps(self, output_path: Optional[str] = None) -> str:
        """Export as PPS (legacy PowerPoint Show, requires LibreOffice)."""
        return self._libreoffice_convert("pps", output_path)

    def to_ppsx(self, output_path: Optional[str] = None) -> str:
        """Export as PPSX (PowerPoint Show, OOXML)."""
        # PPSX is the same container as PPTX — just save with the .ppsx extension.
        output_path = self._output_path("ppsx", output_path)
        self._prs.save(output_path)
        return output_path
